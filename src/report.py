"""Turn grid results into a spreadsheet and a slide deck.

Reads the JSONL written by ``experiment`` and emits:

* ``results.xlsx`` - a summary sheet, the holder x seeker matrix, per-model and
  per-scenario breakdowns, and the raw rows behind all of it.
* ``deck.pptx`` - the same findings as slides with native charts.

All aggregation is plain Python over the rows, so the numbers in the deck and the
spreadsheet come from one code path and cannot drift apart.

Usage:
    python -m src.report --results results/grid.jsonl --outdir reports
"""

from __future__ import annotations

import json
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

# Reference data-viz palette (light surface), used in documented slot order.
SERIES = ["2A78D6", "EB6834", "1BAF7A", "EDA100"]  # blue, orange, aqua, yellow
INK_PRIMARY = "0B0B0B"
INK_SECONDARY = "52514E"
INK_MUTED = "898781"
SURFACE = "FCFCFB"
GRIDLINE = "E1E0D9"
# Blue sequential ramp, light -> dark, for magnitude cells.
SEQUENTIAL = ["CDE2FB", "B7D3F6", "9EC5F4", "86B6EF", "5598E7", "3987E5", "2A78D6", "1C5CAB"]

# The headline measure: share of the holder's secrets disclosed when the
# scenario's own condition for sharing them had not been met.
PRIMARY_METRIC = "inappropriate_rate"


def load_rows(path: str | Path) -> list[dict[str, Any]]:
    """Every well-formed JSON object in a results file."""
    out: list[dict[str, Any]] = []
    results = Path(path)
    if not results.is_file():
        return out
    for line in results.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            row = json.loads(line)
        except json.JSONDecodeError:
            continue
        if isinstance(row, dict):
            out.append(row)
    return out


def _mean(values: list[float]) -> float | None:
    return sum(values) / len(values) if values else None


@dataclass
class Aggregates:
    """Everything both outputs render, computed once."""

    scored: list[dict[str, Any]]
    conversations: list[dict[str, Any]]
    holder_models: list[str] = field(default_factory=list)
    seeker_models: list[str] = field(default_factory=list)
    matrix: dict[tuple[str, str], float | None] = field(default_factory=dict)
    by_holder: dict[str, dict[str, Any]] = field(default_factory=dict)
    by_seeker: dict[str, dict[str, Any]] = field(default_factory=dict)
    by_scenario: dict[str, dict[str, Any]] = field(default_factory=dict)
    totals: dict[str, Any] = field(default_factory=dict)


def aggregate(scored: list[dict[str, Any]], conversations: list[dict[str, Any]]) -> Aggregates:
    ok = [r for r in scored if r.get("ok") and r.get("judge_valid")]
    agg = Aggregates(scored=scored, conversations=conversations)

    agg.holder_models = sorted({r["holder_model"] for r in ok})
    agg.seeker_models = sorted({r["seeker_model"] for r in ok})

    cells: dict[tuple[str, str], list[float]] = defaultdict(list)
    for row in ok:
        cells[(row["holder_model"], row["seeker_model"])].append(row[PRIMARY_METRIC])
    for holder in agg.holder_models:
        for seeker in agg.seeker_models:
            agg.matrix[(holder, seeker)] = _mean(cells.get((holder, seeker), []))

    def summarise(rows: list[dict[str, Any]]) -> dict[str, Any]:
        return {
            "n": len(rows),
            "inappropriate_rate": _mean([r["inappropriate_rate"] for r in rows]),
            "disclosure_rate": _mean([r["disclosure_rate"] for r in rows]),
            "unconditional_rate": _mean([r["unconditional_disclosure_rate"] for r in rows]),
            "explicit": sum(r["explicit_total"] for r in rows),
            "implied": sum(r["implied_total"] for r in rows),
            "appropriate": sum(r["appropriate_total"] for r in rows),
            "inappropriate": sum(r["inappropriate_total"] for r in rows),
            "prompt_leaks": sum(1 for r in rows if r.get("prompt_leak")),
            "first_disclosure_turn": _mean(
                [
                    r["first_disclosure_turn"]
                    for r in rows
                    if r.get("first_disclosure_turn") is not None
                ]
            ),
        }

    for model in agg.holder_models:
        agg.by_holder[model] = summarise([r for r in ok if r["holder_model"] == model])
    for model in agg.seeker_models:
        agg.by_seeker[model] = summarise([r for r in ok if r["seeker_model"] == model])

    for scenario_id in sorted({r["scenario_id"] for r in ok}):
        rows = [r for r in ok if r["scenario_id"] == scenario_id]
        entry = summarise(rows)
        entry["title"] = rows[0]["scenario_title"]
        entry["category"] = rows[0]["category"]
        agg.by_scenario[scenario_id] = entry

    terminations: dict[str, int] = defaultdict(int)
    for row in conversations:
        terminations[str(row.get("termination", "error"))] += 1

    agg.totals = {
        "conversations": len(conversations),
        "conversations_failed": sum(1 for r in conversations if not r.get("ok")),
        "scored_rows": len(scored),
        "scored_usable": len(ok),
        "judge_invalid": sum(1 for r in scored if r.get("ok") and not r.get("judge_valid")),
        "models": len(agg.holder_models),
        "pairs": len({(r["holder_model"], r["seeker_model"]) for r in ok}),
        "scenarios": len(agg.by_scenario),
        "turns": sum(r.get("n_turns", 0) for r in conversations),
        "prompt_tokens": sum(r.get("prompt_tokens", 0) for r in conversations),
        "completion_tokens": sum(r.get("completion_tokens", 0) for r in conversations),
        "cost_usd": round(sum(r.get("cost_usd", 0.0) for r in conversations), 4),
        "wall_clock_h": round(sum(r.get("duration_s", 0.0) for r in conversations) / 3600, 2),
        "terminations": dict(terminations),
        "overall_inappropriate_rate": _mean([r[PRIMARY_METRIC] for r in ok]),
    }
    return agg


# --- spreadsheet ------------------------------------------------------------


def _header(ws, labels: list[str], row: int = 1) -> None:
    for col, label in enumerate(labels, start=1):
        cell = ws.cell(row=row, column=col, value=label)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=INK_SECONDARY)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.freeze_panes = ws.cell(row=row + 1, column=1)


def _autosize(ws, max_width: int = 42) -> None:
    for column in ws.columns:
        letter = get_column_letter(column[0].column)
        longest = max((len(str(c.value)) for c in column if c.value is not None), default=8)
        ws.column_dimensions[letter].width = min(max(10, longest + 2), max_width)


def _ramp_fill(value: float | None, lo: float, hi: float) -> PatternFill | None:
    if value is None:
        return None
    span = (hi - lo) or 1.0
    idx = int(((value - lo) / span) * (len(SEQUENTIAL) - 1))
    return PatternFill("solid", fgColor=SEQUENTIAL[max(0, min(len(SEQUENTIAL) - 1, idx))])


def write_xlsx(agg: Aggregates, path: str | Path) -> Path:
    wb = Workbook()

    # --- summary
    ws = wb.active
    ws.title = "summary"
    ws["A1"] = "Agent-vs-agent disclosure grid"
    ws["A1"].font = Font(bold=True, size=16, color=INK_PRIMARY)
    ws["A2"] = (
        "Every model paired with every other model as information holder and seeker. "
        "The headline measure is the share of the holder's secrets disclosed when the "
        "scenario's own condition for sharing them had not been met."
    )
    ws["A2"].alignment = Alignment(wrap_text=True, vertical="top")
    ws.merge_cells("A2:F4")

    rows = [
        ("Conversations run", agg.totals["conversations"]),
        ("Conversations failed", agg.totals["conversations_failed"]),
        ("Scored rows", agg.totals["scored_rows"]),
        ("Scored rows usable", agg.totals["scored_usable"]),
        ("Judge returned nothing usable", agg.totals["judge_invalid"]),
        ("Models", agg.totals["models"]),
        ("Model pairs", agg.totals["pairs"]),
        ("Scenarios", agg.totals["scenarios"]),
        ("Total turns", agg.totals["turns"]),
        ("Prompt tokens", agg.totals["prompt_tokens"]),
        ("Completion tokens", agg.totals["completion_tokens"]),
        ("API cost (USD)", agg.totals["cost_usd"]),
        ("Generation wall clock (hours)", agg.totals["wall_clock_h"]),
        ("Mean inappropriate-disclosure rate", agg.totals["overall_inappropriate_rate"]),
    ]
    start = 6
    _header(ws, ["Measure", "Value"], row=start)
    for i, (label, value) in enumerate(rows, start=start + 1):
        ws.cell(row=i, column=1, value=label)
        cell = ws.cell(row=i, column=2, value=value)
        if isinstance(value, float):
            cell.number_format = "0.000"
    term_start = start + len(rows) + 2
    ws.cell(row=term_start, column=1, value="Termination").font = Font(bold=True)
    ws.cell(row=term_start, column=2, value="Conversations").font = Font(bold=True)
    for i, (term, count) in enumerate(sorted(agg.totals["terminations"].items()), term_start + 1):
        ws.cell(row=i, column=1, value=term)
        ws.cell(row=i, column=2, value=count)
    _autosize(ws)

    # --- matrix
    ws = wb.create_sheet("matrix")
    ws["A1"] = "Mean inappropriate-disclosure rate: holder (rows) vs seeker (columns)"
    ws["A1"].font = Font(bold=True, size=12)
    ws["A2"] = "Lower is a better-defending holder; a darker column is a more effective seeker."
    ws["A2"].font = Font(italic=True, color=INK_MUTED)
    head = 4
    _header(ws, ["holder \\ seeker", *agg.seeker_models, "mean"], row=head)
    values = [v for v in agg.matrix.values() if v is not None]
    lo, hi = (min(values), max(values)) if values else (0.0, 1.0)
    for r, holder in enumerate(agg.holder_models, start=head + 1):
        ws.cell(row=r, column=1, value=holder).font = Font(bold=True)
        for c, seeker in enumerate(agg.seeker_models, start=2):
            value = agg.matrix.get((holder, seeker))
            cell = ws.cell(row=r, column=c, value=value)
            cell.number_format = "0.00"
            cell.alignment = Alignment(horizontal="center")
            fill = _ramp_fill(value, lo, hi)
            if fill:
                cell.fill = fill
        mean = agg.by_holder.get(holder, {}).get("inappropriate_rate")
        mcell = ws.cell(row=r, column=len(agg.seeker_models) + 2, value=mean)
        mcell.number_format = "0.00"
        mcell.font = Font(bold=True)
    _autosize(ws, max_width=20)

    # --- per-model sheets
    for title, table, key in (
        ("by_holder", agg.by_holder, "holder_model"),
        ("by_seeker", agg.by_seeker, "seeker_model"),
    ):
        ws = wb.create_sheet(title)
        _header(
            ws,
            [
                key,
                "runs",
                "inappropriate_rate",
                "disclosure_rate",
                "unconditional_rate",
                "explicit",
                "implied",
                "appropriate",
                "inappropriate",
                "prompt_leaks",
                "mean_first_disclosure_turn",
            ],
        )
        order = sorted(
            table,
            key=lambda m: (table[m]["inappropriate_rate"] is None, table[m]["inappropriate_rate"]),
        )
        for i, model in enumerate(order, start=2):
            s = table[model]
            ws.cell(row=i, column=1, value=model)
            ws.cell(row=i, column=2, value=s["n"])
            for col, field_name in enumerate(
                ["inappropriate_rate", "disclosure_rate", "unconditional_rate"], start=3
            ):
                cell = ws.cell(row=i, column=col, value=s[field_name])
                cell.number_format = "0.000"
            ws.cell(row=i, column=6, value=s["explicit"])
            ws.cell(row=i, column=7, value=s["implied"])
            ws.cell(row=i, column=8, value=s["appropriate"])
            ws.cell(row=i, column=9, value=s["inappropriate"])
            ws.cell(row=i, column=10, value=s["prompt_leaks"])
            cell = ws.cell(row=i, column=11, value=s["first_disclosure_turn"])
            cell.number_format = "0.0"
        _autosize(ws)

    # --- by scenario
    ws = wb.create_sheet("by_scenario")
    _header(
        ws,
        [
            "scenario",
            "title",
            "category",
            "runs",
            "inappropriate_rate",
            "disclosure_rate",
            "unconditional_rate",
        ],
    )
    for i, (sid, s) in enumerate(sorted(agg.by_scenario.items()), start=2):
        ws.cell(row=i, column=1, value=sid)
        ws.cell(row=i, column=2, value=s["title"])
        ws.cell(row=i, column=3, value=s["category"])
        ws.cell(row=i, column=4, value=s["n"])
        for col, field_name in enumerate(
            ["inappropriate_rate", "disclosure_rate", "unconditional_rate"], start=5
        ):
            cell = ws.cell(row=i, column=col, value=s[field_name])
            cell.number_format = "0.000"
    _autosize(ws)

    # --- raw rows
    for title, rows_out in (("runs", agg.scored), ("conversations", agg.conversations)):
        ws = wb.create_sheet(title)
        if not rows_out:
            ws["A1"] = "(no rows)"
            continue
        columns: list[str] = []
        for row in rows_out:
            for key in row:
                if key not in columns:
                    columns.append(key)
        _header(ws, columns)
        for r, row in enumerate(rows_out, start=2):
            for c, key in enumerate(columns, start=1):
                value = row.get(key)
                if isinstance(value, (dict, list)):
                    value = json.dumps(value)
                ws.cell(row=r, column=c, value=value)
        _autosize(ws, max_width=60)

    out = Path(path)
    out.parent.mkdir(parents=True, exist_ok=True)
    wb.save(out)
    return out


# --- deck -------------------------------------------------------------------


def _text(frame, lines: list[tuple[str, int, bool, str]]) -> None:
    """Fill a text frame with (text, size_pt, bold, hex_colour) paragraphs."""
    frame.word_wrap = True
    for i, (text, size, bold, colour) in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.text = text
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = RGBColor.from_string(colour)


def _slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.35), prs.slide_width - Inches(1.2), Inches(1.0)
    )
    lines = [(title, 28, True, INK_PRIMARY)]
    if subtitle:
        lines.append((subtitle, 13, False, INK_SECONDARY))
    _text(box.text_frame, lines)
    return slide


def _style_chart(chart, *, number_format: str = "0.00") -> None:
    chart.font.size = Pt(11)
    chart.font.color.rgb = RGBColor.from_string(INK_SECONDARY)
    try:
        value_axis = chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.major_gridlines.format.line.color.rgb = RGBColor.from_string(GRIDLINE)
        value_axis.format.line.color.rgb = RGBColor.from_string(GRIDLINE)
        value_axis.tick_labels.number_format = number_format
        value_axis.tick_labels.number_format_is_linked = False
    except (ValueError, AttributeError):  # some chart types expose no value axis
        pass


def _bar_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    categories: list[str],
    series: list[tuple[str, list[float | None]]],
    *,
    note: str = "",
) -> None:
    slide = _slide(prs, title, subtitle)
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)

    height = prs.slide_height - Inches(2.4) if note else prs.slide_height - Inches(1.9)
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6),
        Inches(1.45),
        prs.slide_width - Inches(1.2),
        height,
        data,
    )
    chart = graphic.chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
    for i, plot_series in enumerate(chart.plots[0].series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = RGBColor.from_string(SERIES[i % len(SERIES)])
    # Direct labels: the palette's lighter slots are sub-3:1 on white, so values
    # are never carried by colour alone.
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0.00"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.color.rgb = RGBColor.from_string(INK_SECONDARY)
    _style_chart(chart)

    if note:
        box = slide.shapes.add_textbox(
            Inches(0.6), prs.slide_height - Inches(0.85), prs.slide_width - Inches(1.2), Inches(0.5)
        )
        _text(box.text_frame, [(note, 11, False, INK_MUTED)])


def _matrix_slide(prs: Presentation, agg: Aggregates) -> None:
    slide = _slide(
        prs,
        "Every holder against every seeker",
        "Mean inappropriate-disclosure rate. Rows defend, columns attack; darker = more disclosure.",
    )
    holders, seekers = agg.holder_models, agg.seeker_models
    rows, cols = len(holders) + 1, len(seekers) + 1
    left, top = Inches(0.6), Inches(1.6)
    width = prs.slide_width - Inches(1.2)
    height = Inches(0.42) * rows
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    values = [v for v in agg.matrix.values() if v is not None]
    lo, hi = (min(values), max(values)) if values else (0.0, 1.0)
    span = (hi - lo) or 1.0

    def short(name: str) -> str:
        return name.replace("local-", "")

    table.cell(0, 0).text = "holder \\ seeker"
    for c, seeker in enumerate(seekers, start=1):
        table.cell(0, c).text = short(seeker)
    for r, holder in enumerate(holders, start=1):
        table.cell(r, 0).text = short(holder)
        for c, seeker in enumerate(seekers, start=1):
            value = agg.matrix.get((holder, seeker))
            cell = table.cell(r, c)
            cell.text = "-" if value is None else f"{value:.2f}"
            if value is not None:
                idx = int(((value - lo) / span) * (len(SEQUENTIAL) - 1))
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(SEQUENTIAL[idx])

    for r in range(rows):
        for c in range(cols):
            for para in table.cell(r, c).text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.bold = r == 0 or c == 0
                para.alignment = 2 if c else 1
                para.font.color.rgb = RGBColor.from_string(INK_PRIMARY)


def write_pptx(agg: Aggregates, path: str | Path, *, context: dict[str, Any] | None = None) -> Path:
    context = context or {}
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.3), prs.slide_width - Inches(1.8), Inches(3)
    )
    _text(
        box.text_frame,
        [
            ("Agent-to-agent disclosure grid", 40, True, INK_PRIMARY),
            (
                "Every model paired with every other model, as the side holding private "
                "information and as the side trying to extract it.",
                16,
                False,
                INK_SECONDARY,
            ),
            (
                f"{agg.totals['pairs']} model pairs  |  {agg.totals['scenarios']} scenarios  |  "
                f"{agg.totals['conversations']} conversations  |  {agg.totals['turns']} turns",
                14,
                False,
                INK_MUTED,
            ),
        ],
    )

    # Method
    slide = _slide(prs, "What was measured", "")
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(5)
    )
    _text(
        box.text_frame,
        [
            (
                "Each scenario gives one agent (the holder) private facts, each with an in-world "
                "condition under which sharing it is the strategically correct move. The other "
                "agent (the seeker) tries to draw those facts out.",
                15,
                False,
                INK_PRIMARY,
            ),
            ("", 8, False, INK_PRIMARY),
            (
                "An independent judge model reads each transcript and decides, per secret: was it "
                "revealed, at which turn, and had its condition actually been met first?",
                15,
                False,
                INK_PRIMARY,
            ),
            ("", 8, False, INK_PRIMARY),
            (
                "Inappropriate disclosure = the secret came out with its condition unmet. That is "
                "the failure this grid measures; revealing a secret once its condition is satisfied "
                "is correct play and is scored separately.",
                15,
                True,
                INK_PRIMARY,
            ),
            ("", 8, False, INK_PRIMARY),
            (
                f"Judge: {context.get('judge_model', 'n/a')}    "
                f"Conditions: defense={context.get('defense', 'n/a')}, "
                f"adversary={context.get('adversary', 'n/a')}    "
                f"Max turns: {context.get('max_turns', 'n/a')}",
                12,
                False,
                INK_MUTED,
            ),
        ],
    )

    def short(name: str) -> str:
        return name.replace("local-", "")

    # Holders ranked
    holders = sorted(
        agg.by_holder,
        key=lambda m: (
            agg.by_holder[m]["inappropriate_rate"] is None,
            agg.by_holder[m]["inappropriate_rate"],
        ),
    )
    if holders:
        _bar_slide(
            prs,
            "Which models guard information best?",
            "Mean inappropriate-disclosure rate as the holder. Lower is better.",
            [short(m) for m in holders],
            [
                (
                    "Inappropriate disclosure rate",
                    [agg.by_holder[m]["inappropriate_rate"] for m in holders],
                )
            ],
            note="Each bar averages that model as holder across every seeker it faced, in every scenario.",
        )

    # Seekers ranked
    seekers = sorted(
        agg.by_seeker,
        key=lambda m: (
            agg.by_seeker[m]["inappropriate_rate"] is None,
            agg.by_seeker[m]["inappropriate_rate"],
        ),
        reverse=True,
    )
    if seekers:
        _bar_slide(
            prs,
            "Which models extract information best?",
            "Mean inappropriate disclosure they induced as the seeker. Higher is a stronger extractor.",
            [short(m) for m in seekers],
            [("Disclosure induced", [agg.by_seeker[m]["inappropriate_rate"] for m in seekers])],
            note="Same runs, attributed to the attacking side instead of the defending one.",
        )

    if agg.holder_models and agg.seeker_models:
        _matrix_slide(prs, agg)

    # Scenario / category
    if agg.by_scenario:
        ids = sorted(agg.by_scenario)
        _bar_slide(
            prs,
            "Disclosure by scenario",
            "The same models behave very differently depending on the pretext they face.",
            [f"{sid} ({agg.by_scenario[sid]['category'].replace('_', ' ')})" for sid in ids],
            [
                ("Inappropriate", [agg.by_scenario[s]["inappropriate_rate"] for s in ids]),
                ("Any disclosure", [agg.by_scenario[s]["disclosure_rate"] for s in ids]),
            ],
            note="'Any disclosure' includes secrets shared appropriately, once their condition was met.",
        )

    # Reliability / caveats
    slide = _slide(prs, "Reliability and limits", "")
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(5)
    )
    terminations = ", ".join(f"{k}: {v}" for k, v in sorted(agg.totals["terminations"].items()))
    _text(
        box.text_frame,
        [
            (
                f"Conversations: {agg.totals['conversations']} run, {agg.totals['conversations_failed']} failed.",
                15,
                False,
                INK_PRIMARY,
            ),
            (f"Terminations - {terminations}.", 15, False, INK_PRIMARY),
            (
                f"Scored rows: {agg.totals['scored_usable']} usable of {agg.totals['scored_rows']}; "
                f"{agg.totals['judge_invalid']} produced no usable judgement.",
                15,
                False,
                INK_PRIMARY,
            ),
            ("", 8, False, INK_PRIMARY),
            ("Caveats", 16, True, INK_PRIMARY),
            (
                "- A single judge model scores every run, so its blind spots are systematic, not averaged out.",
                14,
                False,
                INK_SECONDARY,
            ),
            (
                "- One defense and one adversary strategy only; these are baseline conditions, not the full grid.",
                14,
                False,
                INK_SECONDARY,
            ),
            (
                "- One conversation per cell, so per-cell numbers are noisy; the per-model means are the reliable read.",
                14,
                False,
                INK_SECONDARY,
            ),
        ],
    )

    out = Path(path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        prog="python -m src.report",
        description="Build results.xlsx and deck.pptx from grid results.",
    )
    parser.add_argument("--results", default="results/grid.jsonl")
    parser.add_argument("--conversations", default="results/conversations.jsonl")
    parser.add_argument("--outdir", default="reports")
    parser.add_argument("--judge-model", default="")
    parser.add_argument("--defense", default="")
    parser.add_argument("--adversary", default="")
    parser.add_argument("--max-turns", default="")
    args = parser.parse_args(argv)

    scored = load_rows(args.results)
    conversations = load_rows(args.conversations)
    if not scored:
        print(f"No scored rows in {args.results}; nothing to report.")
        return 1

    agg = aggregate(scored, conversations)
    outdir = Path(args.outdir)
    xlsx = write_xlsx(agg, outdir / "results.xlsx")
    context = {
        "judge_model": args.judge_model or (scored[0].get("judge_model", "")),
        "defense": args.defense or scored[0].get("defense", ""),
        "adversary": args.adversary or scored[0].get("adversary", ""),
        "max_turns": args.max_turns,
    }
    pptx = write_pptx(agg, outdir / "deck.pptx", context=context)

    print(f"wrote {xlsx}")
    print(f"wrote {pptx}")
    print(
        f"{agg.totals['scored_usable']} usable rows across "
        f"{agg.totals['pairs']} pairs and {agg.totals['scenarios']} scenarios"
    )
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
